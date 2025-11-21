#!/usr/bin/env python3
"""
Universal Document Processor for Editorial AI Training
Handles emails, PDFs, DOCX, PPTX, XLSX, and more
"""

import json
import argparse
from pathlib import Path
import mailbox
import re

# Document processors - import as needed
def process_mbox(file_path):
    """Process .mbox email files"""
    print(f"  Processing mbox: {file_path.name}")
    emails = []
    mbox = mailbox.mbox(str(file_path))
    
    for message in mbox:
        try:
            body = None
            if message.is_multipart():
                for part in message.walk():
                    if part.get_content_type() == 'text/plain':
                        body = part.get_payload(decode=True)
                        if body:
                            body = body.decode('utf-8', errors='ignore')
                        break
            else:
                body = message.get_payload(decode=True)
                if body:
                    body = body.decode('utf-8', errors='ignore')
            
            if body and len(body) > 50:
                emails.append(clean_text(body))
        except:
            continue
    
    return emails

def process_pdf(file_path):
    """Process PDF files"""
    print(f"  Processing PDF: {file_path.name}")
    try:
        import PyPDF2
        texts = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    texts.append(clean_text(text))
        return texts
    except ImportError:
        print("    ⚠ PyPDF2 not installed. Skipping PDF.")
        return []

def process_docx(file_path):
    """Process Word documents"""
    print(f"  Processing DOCX: {file_path.name}")
    try:
        from docx import Document
        doc = Document(file_path)
        texts = [clean_text(p.text) for p in doc.paragraphs if p.text.strip()]
        return texts
    except ImportError:
        print("    ⚠ python-docx not installed. Skipping DOCX.")
        return []

def process_pptx(file_path):
    """Process PowerPoint files"""
    print(f"  Processing PPTX: {file_path.name}")
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if shape.text.strip():
                        texts.append(clean_text(shape.text))
        return texts
    except ImportError:
        print("    ⚠ python-pptx not installed. Skipping PPTX.")
        return []

def process_xlsx(file_path):
    """Process Excel files"""
    print(f"  Processing XLSX: {file_path.name}")
    try:
        import pandas as pd
        df = pd.read_excel(file_path, sheet_name=None)
        texts = []
        for sheet_name, sheet_df in df.items():
            for col in sheet_df.columns:
                for val in sheet_df[col]:
                    if pd.notna(val) and isinstance(val, str) and len(val) > 20:
                        texts.append(clean_text(str(val)))
        return texts
    except ImportError:
        print("    ⚠ pandas/openpyxl not installed. Skipping XLSX.")
        return []

def process_txt(file_path):
    """Process plain text files"""
    print(f"  Processing TXT: {file_path.name}")
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    return [clean_text(text)] if text else []

def clean_text(text):
    """Clean and normalize text"""
    # Remove excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    # Remove email quotes
    lines = text.split('\n')
    clean_lines = [l for l in lines if not l.strip().startswith('>')]
    text = '\n'.join(clean_lines)
    return text.strip()

def create_training_data(texts, source_type):
    """Convert texts to training examples"""
    training_data = []
    
    for text in texts:
        # Split into reasonable chunks
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        for para in paragraphs:
            if 30 < len(para) < 1000:  # Good training length
                # Create different types of training examples
                training_data.append({
                    "instruction": "Write in the user's style",
                    "input": f"Topic from {source_type}",
                    "output": para
                })
                
                # Also create editing examples
                if len(para) > 100:
                    training_data.append({
                        "instruction": "Edit for clarity and concision",
                        "input": para,
                        "output": para  # In real training, this would be an edited version
                    })
    
    return training_data

def main():
    parser = argparse.ArgumentParser(description="Universal document processor")
    parser.add_argument("--input-dir", default="data/to_process/", 
                       help="Directory with files to process")
    parser.add_argument("--output-dir", default="data/processed/",
                       help="Output directory for training data")
    parser.add_argument("--type", choices=['email', 'pdf', 'docx', 'pptx', 'xlsx', 'txt', 'all'],
                       default='all', help="File type to process")
    
    args = parser.parse_args()
    
    input_dir = Path(args.input_dir)
    output_dir = Path(args.output_dir)
    output_dir.mkdir(exist_ok=True, parents=True)
    
    # File type mappings
    processors = {
        '.mbox': ('email', process_mbox),
        '.pdf': ('pdf', process_pdf),
        '.docx': ('docx', process_docx),
        '.pptx': ('pptx', process_pptx),
        '.xlsx': ('xlsx', process_xlsx),
        '.txt': ('txt', process_txt),
    }
    
    all_training_data = []
    file_counts = {}
    
    print(f"Processing files from: {input_dir}")
    print(f"File type filter: {args.type}")
    print("="*50)
    
    # Process files based on type
    for extension, (doc_type, processor) in processors.items():
        if args.type != 'all' and args.type != doc_type:
            continue
            
        files = list(input_dir.glob(f"*{extension}"))
        if files:
            print(f"\nFound {len(files)} {extension} files")
            file_counts[doc_type] = 0
            
            for file_path in files:
                try:
                    texts = processor(file_path)
                    if texts:
                        training_data = create_training_data(texts, doc_type)
                        all_training_data.extend(training_data)
                        file_counts[doc_type] += 1
                        print(f"    ✓ Extracted {len(training_data)} examples")
                except Exception as e:
                    print(f"    ✗ Error: {str(e)[:50]}")
    
    print("\n" + "="*50)
    print("PROCESSING COMPLETE")
    print("="*50)
    
    # Summary
    for doc_type, count in file_counts.items():
        print(f"  {doc_type}: {count} files processed")
    
    print(f"\nTotal training examples: {len(all_training_data)}")
    
    if all_training_data:
        # Split into train/validation
        split_idx = int(len(all_training_data) * 0.9)
        train = all_training_data[:split_idx]
        val = all_training_data[split_idx:]
        
        # Save
        with open(output_dir / "personal_style_train.jsonl", "w") as f:
            for item in train:
                f.write(json.dumps(item) + "\n")
        
        with open(output_dir / "personal_style_validation.jsonl", "w") as f:
            for item in val:
                f.write(json.dumps(item) + "\n")
        
        print(f"\n✓ Saved training data:")
        print(f"  Training: {len(train)} examples")
        print(f"  Validation: {len(val)} examples")
        print(f"  Location: {output_dir}")
    else:
        print("\n⚠ No training data generated. Check your input files.")

if __name__ == "__main__":
    main()
